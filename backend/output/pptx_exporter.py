"""
pptx_exporter.py — Pure Python PPTX export using python-pptx.
No Node.js required.
"""
import io
import re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Patterns that indicate an error/meta slide that should be suppressed
_BAD_PATTERNS = re.compile(
    r'(technical.{0,30}(issue|problem|error|difficult)'
    r'|(could not|unable to|failed to).{0,40}(fetch|retriev|access|load)'
    r'|data.{0,30}(unavailable|not available|missing|not found)'
    r'|(api|service|network).{0,30}(error|fail|timeout)'
    r'|(sorry|apolog).{0,40}(data|info|content)'
    r'|(fetch|retriev).{0,30}(error|fail)'
    r'|error.{0,30}(occur|encounter))',
    re.IGNORECASE
)


def _is_bad_slide(title: str, bullets: list) -> bool:
    combined = (title + ' ' + ' '.join(bullets)).strip()
    return bool(_BAD_PATTERNS.search(combined))


# ── Theme colours ────────────────────────────────────────
PRIMARY   = RGBColor(0x1E, 0x40, 0xAF)   # deep blue
ACCENT    = RGBColor(0x10, 0xB9, 0x81)   # emerald
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DARK      = RGBColor(0x0F, 0x17, 0x2A)
LIGHT_BG  = RGBColor(0xF1, 0xF5, 0xF9)
MUTED     = RGBColor(0x64, 0x74, 0x8B)


def _tf(shape, text, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    """Replace entire text frame with a single styled run."""
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    p.clear()
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _add_bullet(tf, text: str, size=16, color=DARK, indent=0):
    """Append a bullet paragraph to an existing text frame."""
    p = tf.add_paragraph()
    p.level = indent
    run = p.add_run()
    # Strip leading bullet markers
    clean = re.sub(r'^[-*•]\s*', '', text).strip()
    clean = re.sub(r'\*\*(.*?)\*\*', r'\1', clean)
    run.text = clean
    run.font.size = Pt(size)
    run.font.color.rgb = color


def _solid_fill(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def _add_title_slide(prs: Presentation, topic: str):
    """Slide 0: bold title card."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    W, H = prs.slide_width, prs.slide_height

    # Full background
    bg = slide.shapes.add_shape(1, 0, 0, W, H)
    _solid_fill(bg, PRIMARY)
    bg.line.fill.background()

    # Accent bar at bottom
    bar = slide.shapes.add_shape(1, 0, H - Inches(0.5), W, Inches(0.5))
    _solid_fill(bar, ACCENT)
    bar.line.fill.background()

    # Title
    tx = slide.shapes.add_textbox(Inches(1), Inches(2.2), W - Inches(2), Inches(1.5))
    _tf(tx, topic, size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Subtitle
    tx2 = slide.shapes.add_textbox(Inches(1), Inches(3.8), W - Inches(2), Inches(0.6))
    _tf(tx2, "Research Presentation", size=16, bold=False, color=RGBColor(0xBF, 0xDB, 0xFE), align=PP_ALIGN.CENTER)


def _add_content_slide(prs: Presentation, title: str, bullets: list, slide_num: int):
    """Content slide with header stripe and bullet body."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    W, H = prs.slide_width, prs.slide_height

    # Light background
    bg = slide.shapes.add_shape(1, 0, 0, W, H)
    _solid_fill(bg, LIGHT_BG)
    bg.line.fill.background()

    # Header bar
    hdr = slide.shapes.add_shape(1, 0, 0, W, Inches(1.1))
    _solid_fill(hdr, PRIMARY)
    hdr.line.fill.background()

    # Slide number chip
    chip = slide.shapes.add_shape(1, Inches(0.3), Inches(0.22), Inches(0.55), Inches(0.55))
    _solid_fill(chip, ACCENT)
    chip.line.fill.background()
    chip_tx = chip.text_frame.paragraphs[0]
    chip_run = chip_tx.add_run()
    chip_run.text = str(slide_num)
    chip_run.font.size = Pt(13)
    chip_run.font.bold = True
    chip_run.font.color.rgb = WHITE
    chip_tx.alignment = PP_ALIGN.CENTER

    # Title text in header
    tx = slide.shapes.add_textbox(Inches(1.1), Inches(0.15), W - Inches(1.4), Inches(0.8))
    _tf(tx, title, size=20, bold=True, color=WHITE)

    # Body text box
    body = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), W - Inches(1.0), H - Inches(1.6))
    tf = body.text_frame
    tf.word_wrap = True

    first = True
    for bullet in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.level = 0
        run = p.add_run()
        clean = re.sub(r'^[-*•]\s*', '', bullet).strip()
        clean = re.sub(r'\*\*(.*?)\*\*', r'\1', clean)
        run.text = f"▸  {clean}"
        run.font.size = Pt(13)        # slightly smaller so longer bullets fit
        run.font.color.rgb = DARK
        p.space_before = Pt(8)        # breathing room above each bullet
        p.space_after  = Pt(2)        # small gap below


def parse_presentation_to_slides(topic: str, content: str) -> list:
    """Parse LLM slide outline into list of {title, bullets}."""
    slides = []
    current = None

    for line in content.split('\n'):
        s = line.strip()
        if not s:
            continue

        m = re.match(r'^SLIDE\s+\d+[:.]?\s*(.*)', s, re.IGNORECASE)
        if m:
            if current and current.get('title'):
                slides.append(current)
            current = {'title': m.group(1).strip() or f"Slide {len(slides)+1}", 'bullets': []}
            continue

        if s.startswith('---'):
            if current and current.get('title'):
                slides.append(current)
            current = {'title': f"Slide {len(slides)+1}", 'bullets': []}
            continue

        if current is not None:
            hm = re.match(r'^#{1,3}\s+(.*)', s)
            if hm and not current['title']:
                current['title'] = hm.group(1).strip()
                continue
            bm = re.match(r'^[-*•]\s+(.*)', s)
            if bm:
                current['bullets'].append(bm.group(1).strip()[:200])  # allow longer bullets
            elif current.get('title') and len(s) > 5:
                current['bullets'].append(s[:200])

    if current and current.get('title'):
        slides.append(current)

    # Clamp bullets — allow up to 7 per slide
    for sl in slides:
        sl['bullets'] = sl['bullets'][:7]

    # Filter out any slides that slipped through with error/meta content
    slides = [sl for sl in slides if not _is_bad_slide(sl['title'], sl['bullets'])]

    return slides[:15]   # allow up to 15 content slides


def export_to_pptx(topic: str, content: str) -> bytes:
    """Export presentation content to a .pptx file. Returns bytes."""
    slides = parse_presentation_to_slides(topic, content)

    if not slides:
        # Fallback: extract non-empty, non-error sentences from raw content
        lines = [
            l.strip() for l in content.split('\n')
            if l.strip() and not _BAD_PATTERNS.search(l)
            and not re.match(r'^SLIDE\s*\d+', l.strip(), re.IGNORECASE)
            and len(l.strip()) > 10
        ]
        # Split into chunks of 5 bullets each as separate slides
        chunk_size = 5
        if lines:
            chunks = [lines[i:i+chunk_size] for i in range(0, min(len(lines), 50), chunk_size)]
            slides = [{'title': f'{topic} — Part {i+1}', 'bullets': chunk}
                      for i, chunk in enumerate(chunks)]
        else:
            slides = [{'title': topic, 'bullets': ['Research findings on ' + topic]}]

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Title slide
    _add_title_slide(prs, topic)

    # Content slides
    for i, sl in enumerate(slides, 1):
        _add_content_slide(prs, sl['title'], sl['bullets'], i)

    # Thank you slide
    _add_title_slide(prs, "Thank You")

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
